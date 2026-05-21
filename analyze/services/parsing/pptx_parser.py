import asyncio
import io
import re
import zipfile
from xml.etree import ElementTree

from pptx import Presentation

from analyze.services.llm.client import GigaChatClient


DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
PRESENTATION_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
MATH_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"

NS = {
    "a": DRAWING_NS,
    "p": PRESENTATION_NS,
    "m": MATH_NS,
}


MATH_PROPERTY_NODES = {
    "accPr",
    "barPr",
    "boxPr",
    "ctrlPr",
    "dPr",
    "eqArrPr",
    "fPr",
    "funcPr",
    "groupChrPr",
    "limLowPr",
    "limUppPr",
    "mPr",
    "matrixPr",
    "naryPr",
    "phantPr",
    "rPr",
    "radPr",
    "sPrePr",
    "sSubPr",
    "sSubSupPr",
    "sSupPr",
}


def _namespace(tag: str) -> str:
    if tag.startswith("{"):
        return tag[1:].split("}", 1)[0]

    return ""


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _attribute_value(node: ElementTree.Element, attribute_name: str) -> str:
    for name, value in node.attrib.items():
        if _local_name(name) == attribute_name:
            return value

    return ""


def _clean_text(text: str) -> str:
    text = text.replace("\u00a0", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    text = re.sub(r"\s+([,.;:!?])", r"\1", text)

    return text.strip()


def _clean_math_text(text: str) -> str:
    replacements = {
        "\u2212": "-",
        "\u2013": "-",
        "\u2014": "-",
        "\u00d7": "*",
        "\u2219": "*",
        "\u22c5": "*",
        "\u00b7": "*",
    }

    for source, target in replacements.items():
        text = text.replace(source, target)

    return re.sub(r"\s+", " ", text).strip()


class PPTXParser:
    def __init__(self):
        self.gigachat = GigaChatClient()
        self.semaphore = asyncio.Semaphore(3)

    async def parse(self, file_bytes: bytes) -> str:
        presentation = Presentation(io.BytesIO(file_bytes))

        full_text = self._extract_presentation_text(file_bytes)

        image_tasks = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        image_bytes = shape.image.blob

                        image_tasks.append(
                            self._extract_image_text_safe(image_bytes)
                        )
                    except Exception:
                        pass

        if image_tasks:
            image_results = await asyncio.gather(*image_tasks)

            for text in image_results:
                if text:
                    full_text.append(text)

        return "\n".join(full_text)

    def _extract_presentation_text(self, file_bytes: bytes) -> list[str]:
        try:
            return self._extract_presentation_xml_text(file_bytes)
        except Exception as error:
            print(f"PPTX XML text extraction error: {error}")

            return self._fallback_extract_text(file_bytes)

    def _extract_presentation_xml_text(self, file_bytes: bytes) -> list[str]:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as archive:
            slide_files = sorted(
                [
                    name
                    for name in archive.namelist()
                    if re.match(r"ppt/slides/slide\d+\.xml", name)
                ],
                key=lambda value: int(re.search(r"slide(\d+)\.xml", value).group(1)),
            )

            parts = []

            for slide_file in slide_files:
                slide_xml = archive.read(slide_file)

                root = ElementTree.fromstring(slide_xml)

                slide_parts = self._extract_slide_text(root)

                if slide_parts:
                    parts.append("\n".join(slide_parts))

            return parts

    def _extract_slide_text(self, root: ElementTree.Element) -> list[str]:
        parts = []

        for shape in root.findall(".//p:sp", NS):
            text = self._extract_shape_text(shape)

            if text:
                parts.append(text)

        for table in root.findall(".//a:tbl", NS):
            parts.extend(self._extract_table_text(table))

        return parts

    def _extract_shape_text(self, shape: ElementTree.Element) -> str:
        chunks = []

        for paragraph in shape.findall(".//a:p", NS):
            paragraph_chunks = []

            self._walk_text_node(paragraph, paragraph_chunks)

            text = _clean_text("".join(paragraph_chunks))

            if text:
                chunks.append(text)

        return "\n".join(chunks)

    def _extract_table_text(self, table: ElementTree.Element) -> list[str]:
        rows = []

        for row in table.findall(".//a:tr", NS):
            cells = []

            for cell in row.findall("./a:tc", NS):
                paragraphs = []

                for paragraph in cell.findall(".//a:p", NS):
                    paragraph_chunks = []

                    self._walk_text_node(paragraph, paragraph_chunks)

                    text = _clean_text("".join(paragraph_chunks))

                    if text:
                        paragraphs.append(text)

                if paragraphs:
                    cells.append(" ".join(paragraphs))

            if cells:
                rows.append(" | ".join(cells))

        return rows

    def _walk_text_node(
        self,
        node: ElementTree.Element,
        chunks: list[str],
    ) -> None:
        namespace = _namespace(node.tag)
        local_name = _local_name(node.tag)

        if namespace == MATH_NS and local_name in {"oMath", "oMathPara"}:
            math_text = _clean_math_text(self._omml_to_text(node))

            if math_text:
                chunks.append(f" {math_text} ")

            return

        if namespace == DRAWING_NS and local_name == "t":
            chunks.append(node.text or "")
            return

        if namespace == DRAWING_NS and local_name == "br":
            chunks.append("\n")
            return

        for child in node:
            self._walk_text_node(child, chunks)

    def _omml_to_text(self, node: ElementTree.Element) -> str:
        local_name = _local_name(node.tag)

        if local_name == "t":
            return node.text or ""

        if local_name in MATH_PROPERTY_NODES or local_name.endswith("Pr"):
            return ""

        if local_name == "f":
            numerator = self._math_child_text(node, "num")
            denominator = self._math_child_text(node, "den")

            if numerator and denominator:
                return f"({numerator})/({denominator})"

        if local_name in {"sSup", "sup"}:
            base = self._math_child_text(node, "e")
            superscript = self._math_child_text(node, "sup")

            if base and superscript:
                return f"{base}^{superscript}"

        if local_name in {"sSub", "sub"}:
            base = self._math_child_text(node, "e")
            subscript = self._math_child_text(node, "sub")

            if base and subscript:
                return f"{base}_{subscript}"

        if local_name == "sSubSup":
            base = self._math_child_text(node, "e")
            subscript = self._math_child_text(node, "sub")
            superscript = self._math_child_text(node, "sup")

            if base and subscript and superscript:
                return f"{base}_{subscript}^{superscript}"

        if local_name == "rad":
            degree = self._math_child_text(node, "deg")
            expression = self._math_child_text(node, "e")

            if expression:
                return (
                    f"root({degree}, {expression})"
                    if degree
                    else f"sqrt({expression})"
                )

        if local_name == "d":
            expression = self._math_child_text(node, "e")
            start = "("
            end = ")"

            for child in node:
                if _local_name(child.tag) != "dPr":
                    continue

                for property_node in child:
                    property_name = _local_name(property_node.tag)

                    if property_name == "begChr":
                        start = _attribute_value(property_node, "val") or start

                    if property_name == "endChr":
                        end = _attribute_value(property_node, "val") or end

            return f"{start}{expression}{end}" if expression else ""

        if local_name == "nary":
            operator = self._math_child_text(node, "chr")
            lower = self._math_child_text(node, "sub")
            upper = self._math_child_text(node, "sup")
            expression = self._math_child_text(node, "e")

            limits = ""

            if lower or upper:
                limits = f"_{lower}^{upper}"

            return f"{operator}{limits} {expression}".strip()

        return "".join(self._omml_to_text(child) for child in node)

    def _math_child_text(
        self,
        node: ElementTree.Element,
        child_name: str,
    ) -> str:
        for child in node:
            if _local_name(child.tag) == child_name:
                return _clean_math_text(self._omml_to_text(child))

        return ""

    def _fallback_extract_text(self, file_bytes: bytes) -> list[str]:
        presentation = Presentation(io.BytesIO(file_bytes))

        parts = []

        for slide in presentation.slides:
            slide_parts = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = _clean_text(shape.text)

                    if text:
                        slide_parts.append(text)

            if slide_parts:
                parts.append("\n".join(slide_parts))

        return parts

    async def _extract_image_text_safe(self, image_bytes: bytes) -> str:
        async with self.semaphore:
            try:
                return await self.gigachat.extract_text_from_image(image_bytes)
            except Exception as e:
                print(f"PPTX image extraction error: {e}")
                return ""